"""Generate TAMS-branded PowerPoint presentations using the official template."""

import os
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from config import TAMS_PPTX_TEMPLATE
from templates.report_structure import (
    SECTION_TITLES, SECTION_ORDER, CHART_MAP, report_filename,
)

# Brand colors as RGBColor objects
DEEP_BLUE = RGBColor(0x22, 0x2F, 0x62)
LIGHT_BLUE = RGBColor(0x1A, 0x6D, 0xB6)
TURQUOISE = RGBColor(0x6C, 0xB9, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_CARBON = RGBColor(0x0E, 0x1A, 0x24)
SOFT_CARBON = RGBColor(0xB1, 0xB3, 0xB6)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xF8)

# Template layout indices
LAYOUT_TITLE = 0          # "Title Slide" - has CENTER_TITLE + SUBTITLE
LAYOUT_TITLE_ONLY = 1     # "Title Only" - has TITLE placeholder
LAYOUT_BLANK = 2          # "Blank" - branded but no placeholders
LAYOUT_TAM_BG = 3         # "Tam Background" - full branded background
LAYOUT_EMPTY = 7          # "5_Custom Layout" - empty layout
LAYOUT_BLANK_PAGE = 23    # "Balnk Page" - clean blank

_RTL_PATTERN = re.compile(r'[\u0590-\u05FF\u0600-\u06FF\u0750-\u077F\uFB50-\uFDFF\uFE70-\uFEFF]')


def _detect_rtl(text):
    """Detect if text contains Arabic/Hebrew characters requiring RTL layout."""
    return bool(_RTL_PATTERN.search(text))


def _set_paragraph_rtl(paragraph, is_rtl):
    """Set RTL direction on a paragraph via XML properties."""
    if not is_rtl:
        return
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    pPr.set('algn', 'r')


def _delete_all_slides(prs):
    """Delete all existing slides from the presentation, keeping layouts and masters."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def _add_textbox(slide, left, top, width, height, text, font_size=11,
                 font_color=GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    """Add a text box to a slide with optional RTL support."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    is_rtl = _detect_rtl(text)
    if is_rtl:
        alignment = PP_ALIGN.RIGHT
        _set_paragraph_rtl(p, True)

    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def _add_footer(slide, slide_number, date_str, total_slides=None):
    """Add page number, date, and TAM Capital attribution to slide footer."""
    slide_width = Inches(13.33)
    footer_top = Inches(7.05)
    footer_height = Inches(0.35)

    # Thin separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.95),
        slide_width - Inches(1.0), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = SOFT_CARBON
    line.line.fill.background()

    # Left: TAM Capital attribution
    _add_textbox(
        slide, Inches(0.5), footer_top, Inches(4), footer_height,
        "TAM Capital  |  CMA Licensed", font_size=7,
        font_color=SOFT_CARBON, font_name="Arial"
    )

    # Center: Date
    _add_textbox(
        slide, Inches(5.0), footer_top, Inches(3.33), footer_height,
        date_str, font_size=7, font_color=SOFT_CARBON,
        alignment=PP_ALIGN.CENTER, font_name="Arial"
    )

    # Right: Page number
    page_text = f"{slide_number}"
    if total_slides:
        page_text = f"{slide_number} / {total_slides}"
    _add_textbox(
        slide, Inches(9.83), footer_top, Inches(3), footer_height,
        page_text, font_size=7, font_color=SOFT_CARBON,
        alignment=PP_ALIGN.RIGHT, font_name="Arial"
    )


def _add_cover_slide(prs, stock_name, ticker, date_str):
    """Add cover slide using the official template layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # CENTER_TITLE
            ph.text = stock_name.upper()
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
        elif ph.placeholder_format.idx == 1:  # SUBTITLE
            ph.text = f"({ticker})\nComprehensive Investor Report\n{date_str}"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = SOFT_CARBON

    return slide


def _add_toc_slide(prs, sections, date_str):
    """Add a Table of Contents slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "TABLE OF CONTENTS"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = DEEP_BLUE

    # Build TOC with section numbers
    toc_items = []
    num = 1
    for key in SECTION_ORDER:
        if key in sections:
            toc_items.append((num, SECTION_TITLES.get(key, key.replace("_", " ").title())))
            num += 1

    # Create a formatted text box with proper spacing
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (num, title) in enumerate(toc_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Section number
        run_num = p.add_run()
        run_num.text = f"{num:02d}    "
        run_num.font.size = Pt(14)
        run_num.font.bold = True
        run_num.font.color.rgb = LIGHT_BLUE
        run_num.font.name = "Arial"

        # Section title
        run_title = p.add_run()
        run_title.text = title
        run_title.font.size = Pt(14)
        run_title.font.color.rgb = GRAY
        run_title.font.name = "Arial"

    _add_footer(slide, 2, date_str)
    return slide


def _add_section_header_slide(prs, section_number, section_title, date_str, slide_number):
    """Add a section divider/header slide with the section number and title."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TAM_BG])

    # Large section number
    _add_textbox(
        slide, Inches(1.0), Inches(2.5), Inches(3), Inches(1.5),
        f"{section_number:02d}", font_size=60, font_color=WHITE,
        bold=True, font_name="Arial"
    )

    # Section title
    _add_textbox(
        slide, Inches(1.0), Inches(4.0), Inches(10), Inches(1.0),
        section_title, font_size=32, font_color=WHITE,
        bold=True, font_name="Arial"
    )

    # Thin accent line under the number
    line = slide.shapes.add_shape(
        1,  # rectangle
        Inches(1.0), Inches(3.9),
        Inches(2.0), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TURQUOISE
    line.line.fill.background()

    return slide


def _parse_content_to_paragraphs(content, max_lines=16):
    """Parse markdown content into structured paragraphs for slides."""
    lines = content.strip().split("\n")
    paragraphs = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("|"):
            continue  # Skip markdown tables
        if line.startswith("---") or line.startswith("==="):
            continue
        if line.startswith("## "):
            paragraphs.append({"text": line[3:].strip(), "type": "heading", "level": 2})
        elif line.startswith("### "):
            paragraphs.append({"text": line[4:].strip(), "type": "heading", "level": 3})
        elif line.startswith("#### "):
            paragraphs.append({"text": line[5:].strip(), "type": "heading", "level": 4})
        elif line.startswith("- ") or line.startswith("* "):
            text = line[2:].replace("**", "").strip()
            paragraphs.append({"text": text, "type": "bullet"})
        elif re.match(r'^\d+\.\s+', line):
            text = re.sub(r'^\d+\.\s+', '', line).replace("**", "").strip()
            paragraphs.append({"text": text, "type": "numbered"})
        else:
            paragraphs.append({"text": line.replace("**", "").strip(), "type": "body"})

    return paragraphs


def _add_content_slide(prs, title, content, chart_path=None, date_str="",
                       slide_number=0):
    """Add a content slide with title, body text, and optional chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    # Set title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = DEEP_BLUE

    # Content area dimensions (13.33" wide slide)
    content_left = Inches(0.8)
    content_top = Inches(1.8)
    content_width = Inches(11.5)
    content_height = Inches(4.8)

    has_chart = False
    if chart_path and os.path.exists(chart_path):
        try:
            slide.shapes.add_picture(
                chart_path, Inches(6.8), Inches(1.8), width=Inches(5.5)
            )
            content_width = Inches(5.5)
            has_chart = True
        except Exception:
            pass

    # Parse and render content
    paragraphs = _parse_content_to_paragraphs(content, max_lines=16)
    if not paragraphs:
        _add_footer(slide, slide_number, date_str)
        return slide

    txBox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para in enumerate(paragraphs[:16]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        is_rtl = _detect_rtl(para["text"])
        if is_rtl:
            _set_paragraph_rtl(p, True)

        p.space_before = Pt(2)
        p.space_after = Pt(2)

        if para["type"] == "heading":
            run = p.add_run()
            run.text = para["text"]
            if para["level"] == 2:
                run.font.size = Pt(14)
                run.font.color.rgb = DEEP_BLUE
            elif para["level"] == 3:
                run.font.size = Pt(12)
                run.font.color.rgb = LIGHT_BLUE
            else:
                run.font.size = Pt(11)
                run.font.color.rgb = DEEP_BLUE
            run.font.bold = True
            run.font.name = "Arial"
            p.space_before = Pt(8)
        elif para["type"] == "bullet":
            p.level = 0
            run = p.add_run()
            run.text = f"  \u2022  {para['text']}"
            run.font.size = Pt(10)
            run.font.color.rgb = GRAY
            run.font.name = "Arial"
        elif para["type"] == "numbered":
            run = p.add_run()
            run.text = f"     {para['text']}"
            run.font.size = Pt(10)
            run.font.color.rgb = GRAY
            run.font.name = "Arial"
        else:
            run = p.add_run()
            run.text = para["text"]
            run.font.size = Pt(10)
            run.font.color.rgb = GRAY
            run.font.name = "Arial"

    _add_footer(slide, slide_number, date_str)
    return slide


def _add_table_slide(prs, title, headers, rows, date_str="", slide_number=0):
    """Add a slide with a formatted data table.

    Args:
        prs: Presentation object
        title: Slide title
        headers: List of column header strings
        rows: List of lists (each inner list is a row of cell values)
        date_str: Date string for footer
        slide_number: Page number for footer
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = DEEP_BLUE

    num_rows = min(len(rows) + 1, 15)  # Cap at 14 data rows + header
    num_cols = len(headers)

    table_left = Inches(0.8)
    table_top = Inches(1.8)
    table_width = Inches(11.5)
    table_height = Inches(4.8)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Distribute column widths evenly
    col_width = int(table_width / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width

    # Header row styling
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_BLUE
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Arial"

    # Data rows
    for row_idx, row_data in enumerate(rows[:num_rows - 1]):
        for col_idx, value in enumerate(row_data[:num_cols]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value) if value is not None else ""
            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = GRAY
                    run.font.name = "Arial"

    _add_footer(slide, slide_number, date_str)
    return slide


def _add_sources_slide(prs, sources_text, date_str="", slide_number=0):
    """Add a Sources & References slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "SOURCES & REFERENCES"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = DEEP_BLUE

    _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8),
                 sources_text, font_size=9, font_color=GRAY)

    _add_footer(slide, slide_number, date_str)
    return slide


def _add_disclaimer_slide(prs, date_str="", slide_number=0):
    """Add a disclaimer slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Disclaimer"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
        elif ph.placeholder_format.idx == 1:
            ph.text = (
                "This document is for informational purposes only and does not constitute "
                "investment advice. Past performance does not guarantee future results. "
                "TAM Capital and its affiliates may hold positions in the securities discussed. "
                "TAM Capital is regulated by the Capital Market Authority (CMA) of Saudi Arabia.\n\n"
                "All data sourced from public filings, market data providers, and proprietary analysis. "
                "Redistribution without prior written consent is prohibited."
            )
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = SOFT_CARBON

    return slide


def _set_presentation_metadata(prs, stock_name, ticker):
    """Set document metadata properties."""
    prs.core_properties.author = "TAM Capital"
    prs.core_properties.title = f"{stock_name} ({ticker}) - Investor Report"
    prs.core_properties.subject = "Investment Research Report"
    prs.core_properties.comments = (
        f"Generated by TAM Capital Research Agent on "
        f"{datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )
    prs.core_properties.category = "Financial Research"


def generate_pptx_report(stock_name: str, ticker: str, sections: dict,
                         charts: dict = None, output_dir: str = "output",
                         sources=None) -> str:
    """Generate a TAMS-branded PPTX presentation using the official template.

    Args:
        stock_name: Company name
        ticker: Stock ticker
        sections: Dict of section_name -> content text
        charts: Dict of chart_name -> chart file path
        output_dir: Output directory
        sources: SourceCollector instance or formatted string

    Returns:
        Path to generated PPTX file
    """
    os.makedirs(output_dir, exist_ok=True)

    # Open official template
    if os.path.exists(TAMS_PPTX_TEMPLATE):
        prs = Presentation(TAMS_PPTX_TEMPLATE)
        _delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    date_str = datetime.now().strftime("%B %d, %Y")

    # Set document metadata
    _set_presentation_metadata(prs, stock_name, ticker)

    # --- Pre-calculate total slides for footer numbering ---
    # Count: cover(1) + TOC(1) + section_headers + content_slides + sources?(1) + disclaimer(1)
    content_slide_count = 0
    section_count = 0
    for section_key in SECTION_ORDER:
        if section_key not in sections:
            continue
        section_count += 1
        content = sections[section_key]
        lines = content.strip().split("\n")
        chunk_count = max(1, (len(lines) + 15) // 16)
        content_slide_count += chunk_count

    has_sources = bool(sources)
    # total = cover + TOC + (section headers + content slides) + sources + disclaimer
    total_slides = 1 + 1 + section_count + content_slide_count + (1 if has_sources else 0) + 1

    # --- Build slides ---
    current_slide = 1

    # 1. Cover slide (no footer on cover)
    _add_cover_slide(prs, stock_name, ticker, date_str)
    current_slide += 1

    # 2. Table of Contents
    _add_toc_slide(prs, sections, date_str)
    current_slide += 1

    # 3. Section slides
    section_num = 0
    for section_key in SECTION_ORDER:
        if section_key not in sections:
            continue

        section_num += 1
        title = SECTION_TITLES.get(section_key, section_key.replace("_", " ").title())
        content = sections[section_key]

        # Section header/divider slide
        _add_section_header_slide(prs, section_num, title, date_str, current_slide)
        current_slide += 1

        # Chart for this section
        chart_path = None
        if charts:
            chart_key = CHART_MAP.get(section_key)
            if chart_key:
                chart_path = charts.get(chart_key)

        # Split long content into multiple slides
        lines = content.strip().split("\n")
        chunk_size = 16
        for chunk_start in range(0, len(lines), chunk_size):
            chunk = "\n".join(lines[chunk_start:chunk_start + chunk_size])
            slide_title = title if chunk_start == 0 else f"{title} (cont.)"
            _add_content_slide(
                prs, slide_title, chunk,
                chart_path=chart_path if chunk_start == 0 else None,
                date_str=date_str,
                slide_number=current_slide,
            )
            current_slide += 1

    # 4. Sources slide
    if sources:
        if hasattr(sources, 'format_for_pptx'):
            sources_text = sources.format_for_pptx()
        else:
            sources_text = str(sources)
        _add_sources_slide(prs, sources_text, date_str, current_slide)
        current_slide += 1

    # 5. Disclaimer slide (uses LAYOUT_TITLE, no custom footer needed)
    _add_disclaimer_slide(prs, date_str, current_slide)

    # Save
    filename = report_filename(stock_name, "pptx")
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return filepath
