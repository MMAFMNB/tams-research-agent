"""Tests for PPTX, PDF, and DOCX report generators."""

import os
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document


class TestPPTXGenerator:
    def test_generates_valid_pptx(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Saudi Aramco", "2222.SR", sample_sections, output_dir=output_dir
        )
        assert os.path.exists(path)
        assert path.endswith(".pptx")

    def test_pptx_has_cover_and_content_slides(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Test Corp", "TEST", sample_sections, output_dir=output_dir
        )
        prs = Presentation(path)
        # Cover + TOC + (3 sections * (header + content)) + disclaimer = at least 9
        assert len(prs.slides) >= 9

    def test_pptx_cover_has_company_name(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "SABIC", "2010.SR", sample_sections, output_dir=output_dir
        )
        prs = Presentation(path)
        cover = prs.slides[0]
        texts = [sh.text for sh in cover.shapes if hasattr(sh, "text")]
        assert any("SABIC" in t for t in texts)

    def test_pptx_has_section_headers(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        prs = Presentation(path)
        all_text = " ".join(
            sh.text for slide in prs.slides for sh in slide.shapes if hasattr(sh, "text")
        )
        assert "Executive Summary" in all_text
        assert "Fundamental Analysis" in all_text

    def test_pptx_has_footer_elements(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        prs = Presentation(path)
        # Check a content slide (not cover) has TAM Capital footer text
        content_slide = prs.slides[1]  # TOC slide
        texts = [sh.text for sh in content_slide.shapes if hasattr(sh, "text")]
        assert any("TAM Capital" in t for t in texts)

    def test_pptx_metadata(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Aramco", "2222.SR", sample_sections, output_dir=output_dir
        )
        prs = Presentation(path)
        assert prs.core_properties.author == "TAM Capital"
        assert "Aramco" in prs.core_properties.title

    def test_pptx_file_size_reasonable(self, output_dir, sample_sections):
        from generators.pptx_generator import generate_pptx_report

        path = generate_pptx_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        size_mb = os.path.getsize(path) / (1024 * 1024)
        assert size_mb < 10, f"PPTX too large: {size_mb:.1f}MB"


class TestPDFGenerator:
    def test_generates_valid_pdf(self, output_dir, sample_sections):
        from generators.pdf_generator import generate_pdf_report

        path = generate_pdf_report(
            "Saudi Aramco", "2222.SR", sample_sections, output_dir=output_dir
        )
        assert os.path.exists(path)
        assert path.endswith(".pdf")

    def test_pdf_has_content(self, output_dir, sample_sections):
        from generators.pdf_generator import generate_pdf_report

        path = generate_pdf_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        size = os.path.getsize(path)
        assert size > 5000, f"PDF too small ({size} bytes), likely empty"

    def test_pdf_file_size_reasonable(self, output_dir, sample_sections):
        from generators.pdf_generator import generate_pdf_report

        path = generate_pdf_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        size_mb = os.path.getsize(path) / (1024 * 1024)
        assert size_mb < 5, f"PDF too large: {size_mb:.1f}MB"

    def test_pdf_with_tables_in_content(self, output_dir, sample_sections):
        from generators.pdf_generator import generate_pdf_report

        path = generate_pdf_report(
            "Table Test", "TEST", sample_sections, output_dir=output_dir
        )
        assert os.path.exists(path)
        assert os.path.getsize(path) > 5000


class TestDOCXGenerator:
    def test_generates_valid_docx(self, output_dir, sample_sections):
        from generators.docx_generator import generate_docx_report

        path = generate_docx_report(
            "Saudi Aramco", "2222.SR", sample_sections, output_dir=output_dir
        )
        assert os.path.exists(path)
        assert path.endswith(".docx")

    def test_docx_has_cover_page(self, output_dir, sample_sections):
        from generators.docx_generator import generate_docx_report

        path = generate_docx_report(
            "SABIC", "2010.SR", sample_sections, output_dir=output_dir
        )
        doc = Document(path)
        first_texts = " ".join(p.text for p in doc.paragraphs[:10])
        assert "SABIC" in first_texts

    def test_docx_has_section_headings(self, output_dir, sample_sections):
        from generators.docx_generator import generate_docx_report

        path = generate_docx_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        doc = Document(path)
        all_text = " ".join(p.text for p in doc.paragraphs)
        assert "EXECUTIVE SUMMARY" in all_text
        assert "FUNDAMENTAL ANALYSIS" in all_text

    def test_docx_has_disclaimer(self, output_dir, sample_sections):
        from generators.docx_generator import generate_docx_report

        path = generate_docx_report(
            "Test", "TEST", sample_sections, output_dir=output_dir
        )
        doc = Document(path)
        all_text = " ".join(p.text for p in doc.paragraphs)
        assert "DISCLAIMER" in all_text
        assert "CMA" in all_text or "Capital Market Authority" in all_text
